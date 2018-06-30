# -*- coding: utf-8 -*-

from numpy import average, transpose, zeros, array, linspace, mean, cov, convolve, arange, cumsum, concatenate, abs, \
    mod, std, ones, diff, interp
from pandas import read_csv, read_excel, Timedelta, DataFrame, concat, ols, Series
from scipy import interpolate
import matplotlib.pyplot as plt
from re import match
from pickle import dump, load
from mdfreader import mdfreader
from pptx import Presentation
from pptx.util import Inches
from datetime import date
import time
from requests import get
from math import sin, cos, sqrt, log


class SystemGain(object):
    """
    Main class of system gain, contains all the thing needed to be calculated or plotted.

    Created By Lvhuijia 20180202
    Edited By Luchao 20180204
    Edited By ShenLiang 20180313

    Contains：
    ******
    fun sg_main————Main function of calculating system gain, save  'self.sysGain_class'  to be called from UI.
    ******
    fun acc_response————Prepare acceleration response fig data
    fun launch————Prepare launch fig data
    fun max_acc————Prepare maximum acceleration fig data
    fun pedal_map————Prepare pedal map fig data
    fun shift_map————Prepare shift map fig data
    stc fun cut_sg_data_pedal————Data cutting function
    stc fun arm_interpolate————2-D interpolation method to generate system gain arm fig
    ******
    class AccResponse/Launch/MaxAcc/PedalMap/ShiftMap————Wrapper class of the corresponding fig
    class SystemGainDocker————A class that used to wrap all raw data and results
    ******
    fun plot_max_acc————Plotting method of generating maximum acceleration fig. NEED TO BE REWRITE IN Generate_Figs.py
    fun plot_pedal_map————Plotting method of generating pedal map fig. NEED TO BE REWRITE IN Generate_Figs.py
    fun plot_shift_map————Plotting method of generating shift map fig. NEED TO BE REWRITE IN Generate_Figs.py
    ****** SL
    Add constSpd IN class AccResponse————Wrapper class of the corresponding fig
    Add MaxAcc and conSpd in plot_acc_response_
    ******
    """

    def __init__(self, raw_data, feature_array, pt_type, frequency=20,
                 data_cut_time_of_creep=20, data_cut_time_of_pedal=20,
                 **kwargs):
        """
        Initial function of system gain.

        :param
        """
        self.raw_data = raw_data
        self.feature_array = feature_array
        self.pt_type = pt_type
        self.sysGain_class = {}
        self.frequency = frequency
        self.data_cut_time_of_creep = data_cut_time_of_creep
        self.data_cut_time_of_pedal = data_cut_time_of_pedal
        if 'replace' in kwargs:
            self.replace_cs = True
            self.cs_table = kwargs['cs_table']
        else:
            self.replace_cs = False

    def sg_main(self):
        '''
        Main function of calculating system gain, save  'self.sysGain_class'  to be called from UI.

        '''

        # *******1-GetSysGainData******
        # 获取数据，判断数据类型，不同读取，获取文件名信息，
        # SG_csv_Data_ful = read_csv(file_path)
        self.SG_csv_Data_ful = self.raw_data
        # *******2-GetSGColumn*********
        # 获取列号，标准变量及面板输入，数据预处理
        self.SG_csv_Data_Selc = self.SG_csv_Data_ful.loc[:, self.feature_array]
        # SG_csv_Data = SG_csv_Data_Selc.drop_duplicates()    # 先不去重
        SG_csv_Data = self.SG_csv_Data_Selc  # 选到重复的会报错 ！！！！！所以修改了以下索引格式
        self.time_Data = SG_csv_Data.iloc[:, 0].tolist()
        self.vehSpd_Data = SG_csv_Data.iloc[:, 1].tolist()
        self.pedal_Data = SG_csv_Data.iloc[:, 2].tolist()
        self.acc_Data = SG_csv_Data.iloc[:, 3].tolist()
        self.gear_Data = SG_csv_Data.iloc[:, 4].tolist()
        self.torq_Data = SG_csv_Data.iloc[:, 5].tolist()
        self.fuel_Data = SG_csv_Data.iloc[:, 6].tolist()
        self.enSpd_Data = SG_csv_Data.iloc[:, 7].tolist()
        self.turbSpd_Data = SG_csv_Data.iloc[:, 8].tolist()

        self.colour_Bar = ['orange', 'lightgreen', 'c', 'royalblue', 'lightcoral', 'yellow', 'red', 'brown',
                           'teal', 'blue', 'coral', 'gold', 'lime', 'olive']
        # 数据检查
        if isinstance(self.time_Data[0], int) or isinstance(self.time_Data[0], float):  # 如果时间序列为数字格式
            if sum([int(i <= 0) for i in diff(self.time_Data)]) > 0:  # 如果时间序列是非严格递增的，替换为人工生成的序列
                self.time_Data = arange(0, len(self.time_Data) / self.frequency, 1 / self.frequency)
        else:  # 如果时间序列为非数字格式
            self.time_Data = arange(0, len(self.time_Data) / self.frequency, 1 / self.frequency)

        if not isinstance(self.gear_Data[0], int):  # 如果gear不是int类型的，尝试转换
            try:
                self.gear_Data = [int(i) for i in self.gear_Data]
            except ValueError:
                pass

        self.acc_offset = ((self.vehSpd_Data[-1] - self.vehSpd_Data[0]) / 3.6 / (self.time_Data[-1] - self.time_Data[0])
                           - average(self.acc_Data[:]) * 9.8) / 9.8
        self.acc_Data = [x + self.acc_offset for x in self.acc_Data]

        self.pedal_cut_index, self.pedal_avg = self.cut_sg_data_pedal(self.pedal_Data, self.vehSpd_Data, self.frequency,
                                                                      self.data_cut_time_of_creep,
                                                                      self.data_cut_time_of_pedal)

        trigger = False
        for i, item in enumerate(self.pedal_avg):
            for j, item_ in enumerate(self.pedal_avg):
                if i != j and abs(item - item_) < 3:
                    trigger = True
                    break

        if trigger:  # 挂起线程，等待eliminate_duplicate_ped()被触发
            return 'eliminate_duplicate_ped'
        else:
            self.sg_obj_cal()

    def eliminate_duplicate_ped(self, remove_ped_list):
        ped_h_index = []
        ped_t_index = []
        ped_re = []
        for i, checked in enumerate(remove_ped_list):
            if not checked:
                ped_h_index.append(self.pedal_cut_index[0][i])
                ped_t_index.append(self.pedal_cut_index[1][i])
                ped_re.append(self.pedal_avg[i])
        self.pedal_cut_index = [ped_h_index, ped_t_index]
        self.pedal_avg = ped_re
        self.sg_obj_cal()

    def sg_obj_cal(self):

        self.sort_ped()

        obj_Launch = self.launch(self.acc_Data, self.pedal_Data, self.pedal_cut_index, self.pedal_avg)

        obj_MaxAcc = self.max_acc(self.acc_Data, self.pedal_cut_index, self.pedal_avg)

        obj_PedalMap = self.pedal_map(self.pedal_Data, self.enSpd_Data, self.torq_Data, self.pedal_cut_index,
                                      self.pedal_avg, self.colour_Bar)

        obj_ShiftMap = self.shift_map(self.pedal_Data, self.gear_Data, self.vehSpd_Data, self.pedal_cut_index,
                                      self.pedal_avg, self.colour_Bar,
                                      self.enSpd_Data, self.turbSpd_Data, kind=self.pt_type)
        # # 加速响应曲面增加加速度线及稳车速线
        obj_AccResponse = self.acc_response(self.vehSpd_Data, self.acc_Data, self.pedal_cut_index, self.pedal_avg)

        obj_SystemGain = self.systemgain_curve(self.pedal_Data, self.vehSpd_Data, self.acc_Data, self.pedal_cut_index,
                                               self.pedal_avg)
        self.sysGain_class = self.SystemGainDocker(self.pt_type, obj_AccResponse, obj_Launch, obj_MaxAcc, obj_PedalMap,
                                                   obj_ShiftMap, self.SG_csv_Data_ful, obj_SystemGain)
        return

    def acc_response(self, vehspd_data, acc_data, pedal_cut_index, pedal_avg):
        try:
            acc_ped_map = [[], [], []]
            ped_maxacc = [[], [], []]
            for i in range(0, len(pedal_avg)):
                ivehspd = vehspd_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]
                iped = [pedal_avg[i] * ix / ix for ix in range(pedal_cut_index[0][i], pedal_cut_index[1][i])]
                iacc = acc_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]

                # 删除加速度小于0.05g的数据，认为是踩了制动
                data_len = len(iacc)
                j = 0
                while j < data_len:
                    if iacc[j] < 0.0:  # 20180211 0.05->0.0
                        del ivehspd[j]
                        del iped[j]
                        del iacc[j]
                        data_len = data_len - 1
                        j = j - 1
                    j = j + 1

                acc_ped_map[0].append(iped)
                acc_ped_map[1].append(ivehspd)
                acc_ped_map[2].append(iacc)

                ped_maxacc[0].append(pedal_avg[i])
                ped_maxacc[1].append(ivehspd[iacc.index(max(iacc))])
                ped_maxacc[2].append(max(iacc))

            obj = self.AccResponse(acc_ped_map, pedal_avg, ped_maxacc)
        except Exception:
            acc_ped_map = [[], [], []]
            ped_maxacc = [[], [], []]
            print('Error From acc_response Cal!')
            obj = self.AccResponse(acc_ped_map, pedal_avg, ped_maxacc)
        return obj

    def launch(self, acc_data, pedal_data, pedal_cut_index, pedal_avg):
        try:
            launch_map = [[], [], [], {}]
            delay_acc = []
            delay_time = []
            peak_acc = []
            peak_time = []
            peak_ratio = []

            pedal_cut_index[0] = [int(max(i - 0.5 * self.frequency, 0)) for i in pedal_cut_index[0]]  # max的作用为防止向前超索引

            for i in range(0, len(pedal_avg)):
                if int(round(pedal_avg[i] / 5) * 5) in [10, 20, 30, 50, 100] or (pedal_avg[i] == max(pedal_avg)):
                    iTime = [1/self.frequency * (ix - pedal_cut_index[0][i]) - 0.5 for ix in
                             range(pedal_cut_index[0][i], pedal_cut_index[1][i])]

                    iAcc = acc_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]
                    launch_map[0].append(pedal_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]])
                    launch_map[1].append(iAcc)
                    launch_map[2].append(iTime)

                    for j in range(0, len(iTime)):
                        if iAcc[j] >= 0.02:
                            delay_acc.append(iAcc[j])
                            delay_time.append(iTime[j])
                            break
                        if j == len(iTime)-1:  # 若到最后都没找到
                            delay_acc.append(999)
                            delay_time.append(999)

                    for j in range(0, len(iTime)):
                        if iAcc[j] >= 0.95 * max(iAcc):
                            peak_acc.append(iAcc[j])
                            peak_time.append(iTime[j])
                            peak_ratio.append(iAcc[j]/iTime[j])
                            break

            launch_map[3] = {'delay_acc': delay_acc, 'delay_time': delay_time,
                             'peak_acc': peak_acc, 'peak_time': peak_time, 'peak_ratio': peak_ratio}

            # elif pedal_avg[i] == max(pedal_avg):
            #     iTime = [0.05 * (ix - pedal_cut_index[0][i]) for ix in
            #              range(pedal_cut_index[0][i], pedal_cut_index[0][i] + 100)]
            #     iAcc = acc_data[pedal_cut_index[0][i]:pedal_cut_index[0][i] + 100]
            #     launch_map[0].append(pedal_data[pedal_cut_index[0][i]:pedal_cut_index[0][i] + 100])
            #     launch_map[1].append(iAcc)
            #     launch_map[2].append(iTime)
            obj = self.Launch(launch_map)
        except Exception:
            launch_map = [[], [], []]
            print('Error From acc_launch Cal!')
            obj = self.Launch(launch_map)
        return obj

    def max_acc(self, acc_data, pedal_cut_index, pedal_avg):
        try:
            acc_Ped_Max = []
            for i in range(0, len(pedal_avg)):
                acc_Ped_Max.append(max(acc_data[pedal_cut_index[0][i]:pedal_cut_index[0][i] + 1000]))
            obj = self.MaxAcc(pedal_avg, acc_Ped_Max)
        except Exception:
            acc_Ped_Max = []
            print('Error From max_acc Cal!')
            obj = self.MaxAcc(pedal_avg, acc_Ped_Max)
        return obj

    def pedal_map(self, pedal_data, enSpd_data, torq_data, pedal_cut_index, pedal_avg, colour):
        try:
            pedal_map = [[], [], []]
            for i in range(0, len(pedal_avg)):
                iTorq = torq_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]
                iEnSpd = enSpd_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]
                pedal_map[0].extend(pedal_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]])
                pedal_map[1].extend(iEnSpd)
                pedal_map[2].extend(iTorq)
            obj = self.PedalMap(pedal_map)
        except Exception:
            pedal_map = [[], [], []]
            print('Error From pedal_map Cal!')
            obj = self.PedalMap(pedal_map)
        return obj

    def shift_map(self, pedal_data, gear_data, vehspd_data, pedal_cut_index, pedal_avg, colour,
                  enSpd_Data, turbSpd_Data, kind):
        try:
            shiftMap = [[], [], []]
            if kind == 'AT/DCT':
                for i in range(1, max(gear_data)):
                    # Gear上升沿下降沿
                    for j in range(1, len(gear_data) - 1):
                        if gear_data[j - 1] == i and gear_data[j] == i + 1:
                            for k in range(0, len(pedal_avg)):
                                if j > pedal_cut_index[0][k] and j < pedal_cut_index[1][k]:
                                    shiftMap[0].append(gear_data[j - 1])
                                    shiftMap[1].append(pedal_data[j - 1])
                                    shiftMap[2].append(vehspd_data[j - 1])
                # 按档位油门车速排序
                shiftMap_Sort = sorted(transpose(shiftMap), key=lambda x: [x[0], x[1], x[2]])
                shiftMap_Data = transpose(shiftMap_Sort)
                obj = self.ShiftMap(shiftMap_Data)
            elif kind == 'CVT':
                for i in range(len(pedal_avg)):
                    shiftMap[0].append(enSpd_Data[pedal_cut_index[0][i]:pedal_cut_index[1][i]])
                    shiftMap[1].append(turbSpd_Data[pedal_cut_index[0][i]:pedal_cut_index[1][i]])
                    shiftMap[2].append(vehspd_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]])
                obj = self.ShiftMap(shiftMap, pedal_avg=pedal_avg)
            elif kind == 'AT':
                pass
        except Exception:
            shiftMap = [[], [], []]
            print('Error From shift_map Cal!')
            obj = self.ShiftMap(shiftMap)
        return obj

    def systemgain_curve(self, pedal_data, vehspd_data, acc_data, pedal_cut_index, pedal_avg):
        try:
            vehspd, ped, acc = [], [], []
            vehspd_sg_for_inter = []
            vehspd_steady_for_inter = []

            for i in range(0, len(pedal_avg)):
                ivehspd = vehspd_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]
                iped = [pedal_avg[i]] * (pedal_cut_index[1][i] - pedal_cut_index[0][i])
                iacc = acc_data[pedal_cut_index[0][i]:pedal_cut_index[1][i]]

                vehspd = vehspd + ivehspd
                ped = ped + iped
                acc = acc + iacc
                vehspd_sg_for_inter.append(max(ivehspd))

            if self.replace_cs:  # 如果有稳态车速实验数据输入
                cs_vspd = self.cs_table['meanspd'].tolist()
                cs_ped = self.cs_table['meanped'].tolist()
                if False:  # 用CS作为外界差值的参考
                    cs_interp = interpolate.interp1d(cs_ped, cs_vspd)
                    for i, content in enumerate(pedal_avg):
                        if min(cs_ped) < content < max(cs_ped):
                            vehspd_steady_for_inter.append(cs_interp(content))
                        else:
                            vehspd_steady_for_inter.append(vehspd_sg_for_inter[i])
                else:   # 直接替换为CS
                    vehspd_steady_for_inter = cs_vspd
                    pedal_avg = cs_ped

            else:
                vehspd_steady_for_inter = vehspd_sg_for_inter

            points_for_inter = zeros((len(acc), 2))
            points_for_inter[:, 0] = vehspd
            points_for_inter[:, 1] = ped
            points_for_inter[:, 1] = points_for_inter[:, 1] / 100 * 51
            acc_for_inter = array(acc)

            pedal_steady_for_inter = [x / 100 * 51 + 12.7 for x in pedal_avg]
            vehspd_steady_tb_inter = linspace(min(vehspd_steady_for_inter),
                                              min(max(vehspd_steady_for_inter), 140),
                                              200)
            pedal_steady_function = interpolate.interp1d(vehspd_steady_for_inter, pedal_steady_for_inter)
            pedal_steady_tb_inter = pedal_steady_function(vehspd_steady_tb_inter)
            points_tb_inter = zeros((len(pedal_steady_tb_inter), 2))
            points_tb_inter[:, 0] = vehspd_steady_tb_inter
            points_tb_inter[:, 1] = pedal_steady_tb_inter

            """
            # steady velocity & pedal checked: same with matlab'
            savetxt('vehspd_steady_for_inter.csv', vehspd_steady_for_inter, delimiter=',')
            savetxt('pedal_steady_for_inter.csv', pedal_steady_for_inter, delimiter=',')
            savetxt('pedal_steady_tb_inter.csv', pedal_steady_tb_inter, delimiter=',')
            savetxt('vehspd_steady_tb_inter.csv', vehspd_steady_tb_inter, delimiter=',')
            """

            grid_z1 = interpolate.griddata(points_for_inter, acc_for_inter, points_tb_inter, method='linear')
            grid_z1 = grid_z1 / 12.7
            grid_z1 = self.smooth(grid_z1, 5)

            """
            # function "griddata" checked: same with matlab'
            savetxt('points_for_inter.csv', points_for_inter, delimiter=',')
            savetxt('acc_for_inter.csv', acc_for_inter, delimiter=',')
            savetxt('points_tb_inter.csv', points_tb_inter, delimiter=',')
            savetxt('grid_z1.csv', points_tb_inter, delimiter=',')

            obj = SystemGainPlot()
            """

            acc_vts_after_inter_fun = interpolate.interp1d(vehspd_steady_tb_inter, grid_z1)
            acc_vts_after_inter = []
            for i in [10, 20, 40, 60, 80, 100, 120, 140]:
                try:
                    acc_vts_after_inter.append(acc_vts_after_inter_fun(i).tolist())
                except ValueError:
                    acc_vts_after_inter.append('nan')

            cs_ped_vts_after_inter_fun = interpolate.interp1d(vehspd_steady_for_inter, pedal_avg)
            cs_ped_vts_after_inter = []
            for i in [10, 20, 40, 60, 80, 100, 120]:
                try:
                    cs_ped_vts_after_inter.append(cs_ped_vts_after_inter_fun(i).tolist())
                except ValueError:
                    cs_ped_vts_after_inter.append('nan')

            vehspd_steady_for_inter_log = [log(i) for i in vehspd_steady_tb_inter]
            # corr_of_vspdlog_and_acc = self.compute_correlation(vehspd_steady_for_inter_log, grid_z1)
            r2_of_vspdlog_and_acc = round(
                ols(y=Series(vehspd_steady_for_inter_log), x=DataFrame({"k": grid_z1}))._r2_raw
                , 3)
            # grid_z1_steady_for_inter_log = [log(i) for i in grid_z1]
            # r2_of_vspdlog_and_acc = round(
            #     ols(y=Series(vehspd_steady_tb_inter), x=DataFrame({"k": grid_z1_steady_for_inter_log}))._r2_raw, 3)

            obj = self.SystemGainCurve(vehspd_steady_tb_inter, grid_z1, vehspd_steady_for_inter, pedal_avg,
                                       acc_vts_after_inter, cs_ped_vts_after_inter, r2_of_vspdlog_and_acc)

        except Exception as e:
            obj = self.SystemGainCurve([], [], [], [], [], [], [])
            print('Error From systemgain_curve Cal!')
            print(e)
        return obj

    @staticmethod
    def cut_sg_data_pedal(pedal_data, vehspd_data, frequency, data_cut_time_of_creep=20, data_cut_time_of_pedal=20):
        # 数据切分
        # edges detection initialize to avoid additional detection of rising edges/trailing edges
        pedal_data[0], pedal_data[-1] = 0, 0
        vehspd_data[0], vehspd_data[-1] = 0, 0
        # end of edges detection initialize

        for i in range(len(vehspd_data)):  # GPS 车速置零   20180213 LuChao
            if vehspd_data[i] < 0.5:
                vehspd_data[i] = 0

        r_edge_pedal, t_edge_pedal = [], []
        r_edge_vehspd, t_edge_vehspd = [], []

        # Pedal,vehspd上升沿下降沿
        for i in range(1, len(pedal_data) - 1):

            if pedal_data[i - 1] == 0 and pedal_data[i] > 0:
                r_edge_pedal.append(i)
            if pedal_data[i + 1] == 0 and pedal_data[i] > 0:
                t_edge_pedal.append(i)
            if vehspd_data[i - 1] == 0 and vehspd_data[i] > 0:
                r_edge_vehspd.append(i)
            if vehspd_data[i + 1] == 0 and vehspd_data[i] > 0:
                t_edge_vehspd.append(i)
        # 判断个数大于1000个，去重未做！！
        pedal_cut_index, pedal_avg = [[], []], []

        # creep auto detection 修改了对Creep的判定方式，逻辑如下：1、先寻找到最长的0-Pedal段 2、判定满足长度以及车速状态的0-Pedal段
        for j in range(0, len(t_edge_vehspd)):

            counter = 0
            counter_max = 0
            start_trigger = True
            index_longest_start_0 = 0
            index_longest_stop_0 = 0
            index_longest_start_0_buffer = 0

            for index, ped_content in enumerate(pedal_data[r_edge_vehspd[j]:t_edge_vehspd[j]]):  # 寻找最长的 0 Pedal段
                if int(ped_content) == 0 and index < t_edge_vehspd[j] - r_edge_vehspd[j] - 1:
                    counter += 1
                    if start_trigger:
                        index_longest_start_0_buffer = index + r_edge_vehspd[j]
                        start_trigger = False
                else:
                    if counter > counter_max:
                        counter_max = counter
                        index_longest_stop_0 = index + r_edge_vehspd[j]
                        index_longest_start_0 = index_longest_start_0_buffer
                        start_trigger = True
                    counter = 0

            if (index_longest_stop_0 - index_longest_start_0 > frequency * data_cut_time_of_creep) and (
                            3.5 < average(vehspd_data[index_longest_start_0:index_longest_stop_0]) < 8.5):
                pedal_cut_index[0].append(index_longest_start_0 - frequency * 2)  # 车速信号相对于加速度信号有迟滞
                pedal_cut_index[1].append(index_longest_stop_0 - frequency * 2)
                pedal_avg.append(mean(pedal_data[index_longest_start_0:index_longest_stop_0]))
                break

        for j in range(0, len(r_edge_pedal)):
            if t_edge_pedal[j] - r_edge_pedal[j] > frequency * data_cut_time_of_pedal:  # 时间长度
                if cov(pedal_data[r_edge_pedal[j] + frequency * 1:t_edge_pedal[j] - frequency * 1]) < 3:
                    pedal_cut_index[0].append(r_edge_pedal[j])
                    pedal_cut_index[1].append(t_edge_pedal[j])
                    pedal_avg.append(mean(pedal_data[r_edge_pedal[j]:t_edge_pedal[j]]))

        return pedal_cut_index, pedal_avg

    @staticmethod
    def compute_correlation(x, y):
        x_bar = mean(x)
        y_bar = mean(y)
        ssr = 0
        var_x = 0
        var_y = 0

        for i in range(0, len(x)):
            diff_xx_bar = x[i] - x_bar
            diff_yy_bar = y[i] - y_bar
            ssr += (diff_xx_bar * diff_yy_bar)
            var_x += diff_xx_bar ** 2
            var_y += diff_yy_bar ** 2
        sst = sqrt(var_x * var_y)
        return ssr / sst

    def sort_ped(self):
        ped_h_index = []
        ped_t_index = []
        origin_ped = self.pedal_avg.copy()  # 浅拷贝
        self.pedal_avg.sort()
        for item in self.pedal_avg:
            ped_h_index.append(self.pedal_cut_index[0][origin_ped.index(item)])
            ped_t_index.append(self.pedal_cut_index[1][origin_ped.index(item)])
        self.pedal_cut_index = [ped_h_index, ped_t_index]

    class AccResponse:
        def __init__(self, matrix, para1, matrix_max_acc):
            self.xdata = matrix[1]
            self.ydata = matrix[0]
            self.zdata = matrix[2]
            self.data = matrix
            self.pedal_avg = para1
            self.max_acc_ped = matrix_max_acc
            # self.ped_const = pedal_const
            # self.velo_const = velocity_const

    class Launch:
        def __init__(self, matrix):
            self.xdata = matrix[2]
            self.ydata = matrix[1]
            self.pedal = matrix[0]
            self.data = matrix

    class MaxAcc:
        def __init__(self, x, y):
            self.xdata = x
            self.ydata = y
            self.data = [x, y]

    class PedalMap:
        def __init__(self, matrix):
            self.xdata = matrix[1]
            self.ydata = matrix[2]
            self.zdata = matrix[0]
            self.data = matrix

    class ShiftMap:
        def __init__(self, matrix, **kwargs):
            self.xdata = matrix[2]
            self.ydata = matrix[1]
            self.gear = matrix[0]
            self.data = matrix
            self.kwargs = kwargs

    class SystemGainCurve:
        def __init__(self, vehspd_sg, acc_sg, vehspd_cs, pedal_cs, acc_vts_sg, cs_ped_vts, r2_of_vspdlog_and_acc):
            self.vehspd_sg = vehspd_sg
            self.acc_sg = acc_sg
            self.vehspd_cs = vehspd_cs
            self.pedal_cs = pedal_cs
            self.acc_vts_sg = acc_vts_sg
            self.cs_ped_vts = cs_ped_vts
            self.r2_of_vspdlog_and_acc = r2_of_vspdlog_and_acc

    class SystemGainDocker:
        def __init__(self, pt_type, accresponce, launch, maxacc, pedalmap, shiftmap, rawdata, systemgain):
            self.pt_type = pt_type
            self.accresponce = accresponce
            self.launch = launch
            self.maxacc = maxacc
            self.pedalmap = pedalmap
            self.shiftmap = shiftmap
            self.rawdata = rawdata
            self.systemgain = systemgain


    @staticmethod
    def smooth(data, window):
        # a: 1-D array containing the data to be smoothed by sliding method
        # WSZ: smoothing window size
        out0 = convolve(data, ones(window, dtype=int), 'valid') / window
        r = arange(1, window - 1, 2)
        start = cumsum(data[:window - 1])[::2] / r
        stop = (cumsum(data[:-window:-1])[::2] / r)[::-1]
        return concatenate((start, out0, stop))


class ReadFile(object):

    def __init__(self, file_path, resample_rate, intest_data, feature_row, data_row):
        """
        Initial function of Read File.

        :param file_path: path of system gain file in local disc
        """
        self.file_path = file_path
        self.resample_rate = 1 / int(resample_rate)
        self.intest_data = intest_data
        self.feature_row = int(feature_row) - 1
        self.data_row = int(data_row) - 1
        self.file_columns_orig = []
        self.file_columns = []
        self.import_data_ful_ = []
        self.mdf_data = []
        self.import_status = 'OK'
        self.import_data_information = {}

        self.start_read()

    def start_read(self):
        # 匹配名字信息
        try:
            file_name_match = match(r'^(.*/)(\w*)_(\w*)_([\w.]*)$', self.file_path)
            self.import_data_information['project_name'] = file_name_match.group(2)
            self.import_data_information['vehicle_name'] = file_name_match.group(3)
        except AttributeError:
            self.import_data_information['project_name'] = 'unknown'
            self.import_data_information['vehicle_name'] = 'unknown'
        try:
            self.import_data_information['date_name'] = match(r'^(.*)([0-9]{8})(.*)$', self.file_path).group(2)
        except AttributeError:
            self.import_data_information['date_name'] = 'unknown'

        # 正式读取数据
        file_type_match = match(r'^(.+)(\.)(\w*)$', self.file_path)
        file_type = file_type_match.group(3)
        file_type = file_type.lower()  # 小写
        if file_type == 'csv':
            try:
                self.import_data_ful_ = read_csv(self.file_path, low_memory=False)  # utf-8默认编码
            except UnicodeDecodeError:
                self.import_data_ful_ = read_csv(self.file_path, encoding='gb18030')  # gb18030中文编码

        elif file_type == 'xls' or file_type == 'xlsx':
            try:
                if self.intest_data:
                    _data_ful = read_excel(self.file_path, header=0, skiprows=range(self.feature_row))  # 删头
                    self.import_data_ful_ = _data_ful[self.data_row - self.feature_row - 1:]  # 去单位
                else:
                    self.import_data_ful_ = read_excel(self.file_path)
            except Exception:
                self.import_status = 'Error'

        elif file_type == 'mdf' or file_type == 'dat':
            try:
                self.mdf_data = mdfreader.mdf(self.file_path)
                self.mdf_data.convertToPandas(self.resample_rate)
                data_ful = self.mdf_data['master_group']
                data_ful['time'] = (data_ful.index - data_ful.index[0]).astype('timedelta64[ms]') / 1000
                self.import_data_ful_ = data_ful
            except AttributeError:
                self.import_status = 'Error'

        self.file_columns_orig = self.import_data_ful_.columns
        back_up_counter = 0
        for i in self.import_data_ful_.columns:  # 去除'[]'号，防止控件命名问题
            try:
                ma = match(r'^([0-9a-zA-Z/:_.%\u4e00-\u9fa5\-]+)(\[?)([0-9a-zA-Z/:_.%\u4e00-\u9fa5\[\]\-]*)$', i)
                self.file_columns.append(ma.group(1))
            except AttributeError:
                back_up_counter += 1
                self.file_columns.append('signal' + str(back_up_counter) + '_with_unknown_char')
            except TypeError:  # 原始数据格式有问题
                self.import_status = 'Error'


class SaveAndLoad(object):
    def __init__(self):
        pass

    @staticmethod
    def store_result(file_path, store_data):
        output_file = open(file_path, 'wb')
        dump(store_data, output_file)
        output_file.close()
        return

    @staticmethod
    def store_result_in_excel(file_path, store_data):
        store_data = DataFrame(store_data)
        store_data.to_csv(file_path)
        return

    @staticmethod
    def reload_result(file_path_list):
        data_reload = []
        for i in file_path_list:
            input_file = open(i, 'rb')
            data_reload.append(load(input_file))
            input_file.close()
        return data_reload

    @staticmethod
    def store_feature_index(store_index):
        output_file = open('./bin/feature_index.pkl', 'wb')
        dump(store_index, output_file)
        output_file.close()
        return

    @staticmethod
    def reload_feature_index():
        input_file = open('./bin/feature_index.pkl', 'rb')
        data_reload = load(input_file)
        input_file.close()
        return data_reload

    @staticmethod
    def search_feature(cols_to_search, feature):
        if not isinstance(cols_to_search, list):
            try:
                cols_to_search = cols_to_search.tolist()
            except AttributeError:
                return
        for i in feature:
            for j in cols_to_search:
                _match_result = match(i, j)
                try:
                    match_result = _match_result.group(0)
                    match_result_index = cols_to_search.index(match_result)
                    return match_result_index
                except AttributeError:
                    pass
        return 0  # 什么都没找到的时候返回0

    # PPT
    @staticmethod
    def get_fig_name(type, rawdata_filepath):
        if type == "initial_ppt":
            list_figs = [['dr_sg_curve', 'System Gain Curve.png'], ['dr_max_acc', 'Max Acc.png'],
                         ['dr_cons_spd', 'Constant Speed.png'], ['dr_acc_curve', 'Acc Response.png'],
                         ['dr_launch', 'Launch.png'], ['dr_shift_map', 'Shift Map.png'],
                         ['dr_pedal_map', 'PedalMap.png']]
            rawdata_name = rawdata_filepath.split("/")[-1].split(".")[0]
            title_ppt = "Drive Quality Report"
        else:
            # type == "analysis_ppt"
            list_figs = [['dr_history_sg_curve', 'System Gain Curve.png'], ['dr_history_max_acc', 'Max Acc.png'],
                         ['dr_history_cons_spd', 'Constant Speed.png'], ['dr_history_acc_curve', 'Acc Response.png'],
                         ['dr_history_launch', 'Launch.png'], ['dr_history_shift_map', 'Shift Map.png'],
                         ['dr_history_pedal_map', 'PedalMap.png']]
            rawdata_name = ""
            title_ppt = "Drive Quality Compare Result"
        return list_figs, rawdata_name, title_ppt

        # PPT

    @staticmethod
    def save_pic_ppt(list_figs, rawdata_name, title_ppt, pic_path):
        # Input PowerPoint Template，default blank。
        try:
            prs = Presentation("SAIC_Temp.pptx")
        except Exception:
            prs = Presentation()

        date_pos = [5, 6.9, 3, 0.4]
        pic_pos = [1.95, 1.2, 5.7, 5.0]
        title_pos = [4, 6.15, 3, 0.3]
        data_pos = [4.8, 6.85, 3, 0.5]

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_ppt
        title.text_frame.paragraphs[0].font.size = Inches(0.6)
        subtitle.text = 'Data: ' + rawdata_name

        textbox = slide.shapes.add_textbox(Inches(date_pos[0]), Inches(date_pos[1]),
                                           Inches(date_pos[2]), Inches(date_pos[3]))
        textbox.text = "Generated on {:%Y-%m-%d}".format(date.today())

        try:
            for i in range(0, len(list_figs)):
                pic_name = list_figs[i][1].split(".")[0]
                graph_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(graph_slide_layout)
                textbox = slide.shapes.add_textbox(Inches(title_pos[0]), Inches(title_pos[1]),
                                                   Inches(title_pos[2]), Inches(title_pos[3]))
                textbox.text = pic_name
                if i > 2 and title_ppt == "Drive Quality Compare Result":
                    pic_pos = [0.12, 1.5, 9.7, 4.4]

                pic = slide.shapes.add_picture(pic_path + pic_name + '.png', Inches(pic_pos[0]),
                                               Inches(pic_pos[1]), Inches(pic_pos[2]), Inches(pic_pos[3]))
                # pic = placeholder.insert_picture(chart[i])
                textbox = slide.shapes.add_textbox(Inches(data_pos[0]), Inches(data_pos[1]),
                                                   Inches(data_pos[2]), Inches(data_pos[3]))
                textbox.text = rawdata_name
                textbox.text_frame.paragraphs[0].font.size = Inches(0.14)
            # Add the end of PPT
            end_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(end_slide_layout)
            textbox = slide.shapes.add_textbox(Inches(title_pos[0] - 0.8), Inches(title_pos[1] - 3),
                                               Inches(title_pos[2]), Inches(title_pos[3]))
            textbox.text = "Thank You"

            textbox.text_frame.paragraphs[0].font.size = Inches(0.8)
            textbox.text_frame.paragraphs[0].font.bold = True
        except FileNotFoundError as err:
            print('From save_pic_ppt')
            print(err)
        return prs


class ConstantSpeed(object):
    def __init__(self, raw_data, feature_list, frequency=20, time_block=15):
        self.raw_data = raw_data
        self.feature_list = feature_list
        self.frequency = frequency
        self.time_block = time_block
        self.gear_consider = True
        self.calculation_error = False

        self.cs_csv_data = self.raw_data.loc[:, feature_list]
        self.veh_spd = self.cs_csv_data.iloc[:, 1].tolist()
        self.pedal = self.cs_csv_data.iloc[:, 2].tolist()
        self.acc = self.cs_csv_data.iloc[:, 3].tolist()
        self.gear = self.cs_csv_data.iloc[:, 4].tolist()
        self.en_spd = self.cs_csv_data.iloc[:, 7].tolist()
        try:
            self.acc_smo = self.smooth(self.acc, 1.5 * self.frequency)
        except ValueError:
            self.calculation_error = True
            return
        self.cs_table = []
        self.cs_main()

    def cs_main(self):
        k = 0
        self.cs_table.append([0, 5, 0, 0, 0, 800, 1, 0])
        #                    [index, meanspd, meanped, meanacc, acc_smooth_rate, meanrpm, gear , spd_std]

        try:
            gear_nums = {}
            for i in self.gear:  # 统计Gear信号中的数据种类
                gear_nums[i] = gear_nums.get(i, 0) + 1
            if gear_nums.keys().__len__() > 15:  # 若数据种类大于15种，认为是非DCT/MT/AT车型的，可能为CVT，也可能是选错了信号
                self.gear_consider = False

            for i in range(0, len(self.veh_spd) - self.frequency * self.time_block, self.frequency):
                if self.veh_std_cal(i) < 0.7 and self.mean_speed_cal(i) > 5 and \
                        (self.gear_std_cal(i) == 0 and self.gear_consider or not self.gear_consider):
                    # 若为DCT/MT/AT车型，则要求稳态车速实验时挡位固定，否则不要求
                    if k == 0:
                        self.cs_table.append([i,
                                              self.mean_speed_cal(i),
                                              self.mean_ped_cal(i),
                                              self.mean_acc_cal(i),
                                              self.acc_smooth_rate_std_cal(i),
                                              self.mean_rpm_cal(i),
                                              self.gear[i],
                                              self.veh_std_cal(i)])
                        k = k + 1
                    elif self.veh_std_cal(i) < self.cs_table[k][7] and abs(
                                    self.mean_speed_cal(i) - self.cs_table[k][1]) < 5:
                        self.cs_table[k] = [i,
                                            self.mean_speed_cal(i),
                                            self.mean_ped_cal(i),
                                            self.mean_acc_cal(i),
                                            self.acc_smooth_rate_std_cal(i),
                                            self.mean_rpm_cal(i),
                                            self.gear[i],
                                            self.veh_std_cal(i)]
                    elif abs(self.mean_speed_cal(i) - self.cs_table[k][1]) >= 5:
                        k = k + 1
                        self.cs_table.append([i,
                                              self.mean_speed_cal(i),
                                              self.mean_ped_cal(i),
                                              self.mean_acc_cal(i),
                                              self.acc_smooth_rate_std_cal(i),
                                              self.mean_rpm_cal(i),
                                              self.gear[i],
                                              self.veh_std_cal(i)])

            self.cs_table = DataFrame(self.cs_table, columns=['index', 'meanspd', 'meanped', 'meanacc',
                                                              'acc_smooth_rate', 'meanrpm', 'gear', 'spd_std'])
            self.cs_table.sort_values(by=['meanspd'], inplace=True)
            # 排序后请用iloc索引
            pass

        except Exception :
            self.calculation_error = True

    @staticmethod
    def smooth(data, window):
        # a: 1-D array containing the data to be smoothed by sliding method
        # WSZ: smoothing window size
        if not isinstance(window, int):
            window = int(window)
        if mod(window, 2) == 0:
            window = window + 1

        out0 = convolve(data, ones(window, dtype=int), 'valid') / window
        r = arange(1, window - 1, 2)
        start = cumsum(data[:window - 1])[::2] / r
        stop = (cumsum(data[:-window:-1])[::2] / r)[::-1]
        return concatenate((start, out0, stop))

    def veh_std_cal(self, i):
        return std(self.veh_spd[i:i + self.frequency * self.time_block])

    def gear_std_cal(self, i):
        return std(self.gear[i:i + self.frequency * self.time_block])

    def acc_smooth_rate_std_cal(self, i):
        return std(self.acc[i:i + self.frequency * self.time_block] -
                   self.acc_smo[i:i + self.frequency * self.time_block])

    def mean_speed_cal(self, i):
        return mean(self.veh_spd[i:i + self.frequency * self.time_block])

    def mean_ped_cal(self, i):
        return mean(self.pedal[i:i + self.frequency * self.time_block])

    def mean_acc_cal(self, i):
        return mean(self.acc[i:i + self.frequency * self.time_block])

    def mean_rpm_cal(self, i):
        return mean(self.en_spd[i:i + self.frequency * self.time_block])


class RealRoadFc(object):
    x_PI = 3.14159265358979324 * 3000.0 / 180.0
    PI = 3.1415926535897932384626
    a = 6378245.0
    ee = 0.00669342162296594323

    def __init__(self, raw_data, feature_array, pt_type, frequency=20, ):
        self.raw_data = raw_data
        self.feature_array = feature_array
        self.pt_type = pt_type
        self.real_road_class = {}
        self.frequency = frequency

    def real_road_main(self):
        real_road_data_select = self.raw_data.loc[:, self.feature_array]
        time_data = real_road_data_select.iloc[:, 0]
        vehicle_speed_data = real_road_data_select.iloc[:, 1]
        pedal_data = real_road_data_select.iloc[:, 2]
        acceleration_x_data = real_road_data_select.iloc[:, 3]
        gear_data = real_road_data_select.iloc[:, 4]
        torque_data = real_road_data_select.iloc[:, 5]
        fuel_data = real_road_data_select.iloc[:, 6]
        engine_speed_data = real_road_data_select.iloc[:, 7]
        turbine_speed_data = real_road_data_select.iloc[:, 8]
        brake_data = real_road_data_select.iloc[:, 9]
        latitude_data = real_road_data_select.iloc[:, 10]
        longitude_data = real_road_data_select.iloc[:, 11]
        steering_angle_data = real_road_data_select.iloc[:, 12]
        acceleration_y_data = real_road_data_select.iloc[:, 13]

        time_index_start = time_data[0]

        # driver_name =
        # travel_date =
        # travel_id =

        # target_driver = target_driver[1000:]   # 切掉冷启动的

        x_acc = self.five_points_avg_acc(vehicle_speed_data.tolist(), frequency=10)
        y_acc = self.five_points_avg(acceleration_y_data.tolist())

        delta_tc_speed = engine_speed_data - turbine_speed_data
        tc_open_data = time_data[(gear_data > 0) & (abs(delta_tc_speed) > 200)]
        tc_close_data = time_data[(gear_data > 0) & (abs(delta_tc_speed) < 50) & vehicle_speed_data > 1]
        # ------------------- Style features ---------------------
        odometer = self.odometer_cal(vehicle_speed_data)
        fuel_cons = sum(fuel_data) / odometer * 100 / 1e6  # L/100km
        fuel_csump_run = fuel_data[vehicle_speed_data > 0]
        fuel_cons_run = sum(fuel_csump_run) / odometer * 100 / 1e6  # L/100km

        sudden_acc_times = self.sudden_acc(x_acc, vehicle_speed_data.tolist())
        sudden_brake_times = self.sudden_brake(x_acc)
        sudden_steering_times = self.sudden_steering(y_acc)
        tip_in_times = self.tip_in(pedal_data.tolist())

        # ----------------- Statistics features ------------------
        # soc_start = target_driver['BMSPackSOC'].tolist()[0]
        # soc_end = target_driver['BMSPackSOC'].tolist()[-1]
        # soc_balance_ratio = soc_balance_time_ratio(target_driver['BMSPackSOC'])
        mean_spd = vehicle_speed_data.mean()
        mean_spd_run = vehicle_speed_data[vehicle_speed_data > 0].mean()
        std_spd = vehicle_speed_data.std()
        mean_str_whl_ang = steering_angle_data.mean()
        std_str_whl_ang = steering_angle_data.std()
        mean_str_whl_ang_pos = steering_angle_data[steering_angle_data > 0].mean()
        mean_str_whl_ang_neg = steering_angle_data[steering_angle_data < 0].mean()
        mean_x_acc = x_acc.mean()
        std_x_acc = x_acc.std()
        mean_x_acc_pos = x_acc[x_acc > 0].mean()
        mean_x_acc_neg = x_acc[x_acc < 0].mean()
        mean_acc_ped = pedal_data.mean()
        std_acc_ped = pedal_data.std()
        idle_time = time_data[(vehicle_speed_data == 0) & (engine_speed_data > 0)].shape[0] / 10
        run_time = time_data[vehicle_speed_data > 0].shape[0] / 10
        total_time = vehicle_speed_data.shape[0] / 10
        fuel_csump_idle = fuel_data[(vehicle_speed_data == 0) & (engine_speed_data > 0)]
        try:
            fuel_cons_idle = 0.1 * sum(fuel_csump_idle) / idle_time / 100
        except ZeroDivisionError:
            fuel_cons_idle = 9999
        start_times = self.start_times_cal(vehicle_speed_data)

        sudden_acc_times = array(sudden_acc_times) / odometer
        sudden_brake_times = array(sudden_brake_times) / odometer
        sudden_steering_times = array(sudden_steering_times) / odometer

        sudden_acc_score = (sudden_acc_times[0] + sudden_acc_times[1] * 3 + sudden_acc_times[2] * 7) / odometer
        sudden_brake_score = (sudden_brake_times[0] + sudden_brake_times[1] * 3 + sudden_brake_times[
            2] * 7) / odometer
        sudden_steering_score = (sudden_steering_times[0] + sudden_steering_times[1] * 3 + sudden_steering_times[
            2] * 7) / odometer

        # ------------------- Plot Data --------------------------
        # fuel_distri = self.spd_distribution_array(vspd, fuel_csump, type='sum')
        # portion_distri = self.spd_distribution_array(vspd, vspd, type='weight')
        Date = date_reverse(int(Date))
        resultarray.append(
            [str(driver_name), int(travel_date), int(time_index_start), brake_skill_score, sudden_acc_score,
             sudden_brake_score, sudden_steering_score, overtake_times / odemeter])
        backup_data.append([int(TravelID), round(sudden_acc_times[0], 3), round(sudden_acc_times[1], 3),
                            round(sudden_acc_times[2], 3), round(sudden_brake_times[0], 3),
                            round(sudden_brake_times[1], 3), round(sudden_brake_times[2], 3),
                            round(sudden_steering_times[0], 3), round(sudden_steering_times[1], 3),
                            round(sudden_steering_times[2], 3),
                            round(brake_skill_score, 3), round(overtake_times / odemeter, 4),
                            round(tip_in_times[0] / odemeter, 4),
                            round(odemeter, 2), round(mean_spd, 1), round(mean_Strg_Whl_Ang, 1),
                            round(mean_Strg_Whl_Ang_pos, 1), round(mean_Strg_Whl_Ang_neg, 1),
                            round(mean_X_acc, 3), round(mean_X_acc_pos, 3), round(mean_X_acc_neg, 3),
                            round(mean_acc_ped, 2), round(std_spd, 2), round(std_Strg_Whl_Ang, 2),
                            round(std_X_acc, 3), round(std_acc_ped, 2), round(over_speed_propotion, 3),
                            round(fuel_cons, 3), round(soc_balance_ratio, 3), round(soc_start, 1),
                            round(soc_end, 1), round(run_time, 1), round(total_time, 1), round(fuel_cons_idle, 3),
                            int(start_times)])

        rq = requests.get(static_map_request_url(Target_Driver[['Longitude', 'Latitude']].values.tolist()))
        file = open('./RoutinePic/' + str(CarName) + '_' + str(DriverID) + '_' + str(Date) + '_' +
                    str(int(Time_index_start)) + '.png', 'wb')
        file.write(rq.content)
        file.close()

    @staticmethod
    def find_3_level_times(series, base=0, a1=1, a2=2, a3=3, delta_t=20):  # Passed Test 20170526
        count_1 = count_2 = count_3 = 0
        up_1 = up_2 = up_3 = 0
        down_1 = down_2 = down_3 = 0
        l = len(series)
        for j in range(l - 1):
            if series[j] == base and series[j + 1] == a1:
                up_1 = j + 1
            if series[j] == a1 and series[j + 1] == a2:
                up_2 = j + 1
            if series[j] == a2 and series[j + 1] == a3:
                up_3 = j + 1
            if series[j] == a1 and series[j + 1] == base:
                down_1 = j + 1
            if series[j] == a2 and series[j + 1] == a1:
                down_2 = j + 1
            if series[j] == a3 and series[j + 1] == a2:
                down_3 = j + 1
            if series[j] > series[j + 1]:
                if down_3 - up_3 < delta_t <= down_2 - up_2:
                    count_2 = count_2 + 1
                    down_1 = down_2 = down_3 = -l * 2
                    up_1 = up_2 = up_3 = l * 2
                if down_3 - up_3 >= delta_t:
                    count_3 = count_3 + 1
                    down_1 = down_2 = down_3 = -l * 2
                    up_1 = up_2 = up_3 = l * 2
                if down_2 - up_2 < delta_t <= down_1 - up_1:
                    count_1 = count_1 + 1
                    down_1 = down_2 = down_3 = -l * 2
                    up_1 = up_2 = up_3 = l * 2
        return count_1, count_2, count_3

    @staticmethod
    def five_points_avg(series):
        l = len(series)
        avg_series = zeros(l)
        avg_series[0] = series[1]
        avg_series[1] = (series[2] + series[1] + series[0]) / 3
        for i in range(2, l - 2, 1):
            avg_series[i] = (series[i + 2] + series[i + 1] + series[i] + series[i - 1] + series[i - 2]) / 5
        avg_series[-2] = (series[-1] + series[-2] + series[-3]) / 3
        avg_series[-1] = series[-1]
        return avg_series

    @staticmethod
    def five_points_avg_acc(series, g=9.8, frequency=10):
        l = len(series)
        avg_acc = zeros(l)
        avg_acc[0] = (series[1] - series[0]) * frequency / 3.6 / g
        avg_acc[1] = (series[2] - series[0]) * frequency / 2 / 3.6 / g
        for i in range(2, l - 2, 1):
            avg_acc[i] = (series[i + 2] - series[i - 2]) * frequency / 4 / 3.6 / g
        avg_acc[-2] = (series[-1] - series[-3]) * frequency / 2 / 3.6 / g
        avg_acc[-1] = (series[-1] - series[-2]) * frequency / 3.6 / g
        return avg_acc

    @staticmethod
    def fuel_cal(series, buffer_size=65535):
        l = len(series)
        fuel_count = 0
        bug_fc = []
        for i in range(1, l, 1):
            if (series[i] - series[i - 1]) < 0:
                if series[i - 1] > (buffer_size - 500):
                    fuel_count = fuel_count + 1
                else:
                    bug_fc.append(series[i - 1])
        fuel_sum = fuel_count * buffer_size + series[-1] - series[0] + sum(bug_fc)
        return fuel_sum / 1e6  # unit L

    @staticmethod
    def signal_times_cal(series):  # 数上升沿个数  Passed Test
        l = len(series)
        times_counter = 0
        for i in range(1, l, 1):
            if (not series[i - 1]) and series[i]:
                times_counter = times_counter + 1
            else:
                continue
        return

    def date_reverse(self, date_in):
        revised_data = []
        if date_in.__class__ == str:
            revised_data = self.normal_time_stamp(int(date_in))
        elif date_in.__class__ == int:
            revised_data = self.normal_time_stamp(date_in)
        return revised_data

    @staticmethod
    def normal_time_stamp(date_in):
        if 170101 < date_in < 180101:  # 年月日
            normal_time_stamp = (int(date_in / 10000) + 2000) * 10000 + (int(
                (date_in - int(date_in / 10000) * 10000) / 100)) * 100 + int(date_in - int(date_in / 100) * 100)

        else:  # 日月年
            normal_time_stamp = (int(date_in - int(date_in / 100) * 100) + 2000) * 10000 + (int(
                (date_in - int(date_in / 10000) * 10000) / 100)) * 100 + int(date_in / 10000)

        return normal_time_stamp

    @staticmethod
    def utc_mktime(DATE, TIME):
        """Returns number of seconds elapsed since epoch
     
        Note that no timezone are taken into consideration.
     
        utc tuple must be: (year, month, day, hour, minute, second)
        """
        if 170101 < DATE < 180101:  # 年月日
            utc_tuple = (int(DATE / 10000) + 2000,
                         int((DATE - int(DATE / 10000) * 10000) / 100),
                         int(DATE - int(DATE / 100) * 100),
                         int(TIME / 10000),
                         int((TIME - int(TIME / 10000) * 10000) / 100),
                         int(TIME - int(TIME / 100) * 100),
                         0, 0, 0)
        else:  # 日月年
            utc_tuple = (int(DATE - int(DATE / 100) * 100) + 2000,
                         int((DATE - int(DATE / 10000) * 10000) / 100),
                         int(DATE / 10000),
                         int(TIME / 10000),
                         int((TIME - int(TIME / 10000) * 10000) / 100),
                         int(TIME - int(TIME / 100) * 100),
                         0, 0, 0)
        return int(time.mktime(utc_tuple))

    def datetime_to_timestamp(self, dt):
        """Converts a datetime object to UTC timestamp"""
        return int(self.utc_mktime(dt.timetuple()))

    # ----------------------  Calculation functions  --------------------------
    @staticmethod
    def odometer_cal(vspd, frequency=10):
        odometer = sum(vspd) / frequency / 3.6 / 1000  # km
        return odometer

    def sudden_brake(self, series, car_type='default'):
        style_brk = {'ZS11': [-0.5, -0.35, -0.1],
                     'AS24': [-0.5, -0.35, -0.15],
                     'default': [-0.5, -0.35, -0.15]}
        l = len(series)
        brk_lab = zeros(l)
        for i in range(l):
            if series[i] < style_brk[car_type][0]:
                brk_lab[i] = 3
            elif style_brk[car_type][0] < series[i] < style_brk[car_type][1]:
                brk_lab[i] = 2
            elif style_brk[car_type][1] < series[i] < style_brk[car_type][2]:
                brk_lab[i] = 1
        brk_times = self.find_3_level_times(brk_lab)
        return brk_times

    def sudden_acc(self, series_xacc, series_vspd, car_type='default'):
        style_acc = {'ZS11': [0.04, 0.31, 0.36, 0.35, 0.22, 0.23, 0.16, 0.12, 0.08],
                     'AS24': [0.2, 0.38, 0.41, 0.41, 0.39, 0.43, 0.36, 0.30, 0.22],
                     'default': [0.2, 0.38, 0.41, 0.41, 0.39, 0.43, 0.36, 0.30, 0.22]}
        vspd_interp = [0, 10, 20, 30, 40, 60, 80, 100, 120]
        l = len(series_xacc)
        acc_lab = zeros(l)
        style_acc_current = array(style_acc[car_type])
        for i in range(l):
            if series_xacc[i] > interp(series_vspd[i], vspd_interp, 0.5 * style_acc_current):
                acc_lab[i] = 3
            elif interp(series_vspd[i], vspd_interp, 0.3 * style_acc_current) < series_xacc[i] < \
                    interp(series_vspd[i], vspd_interp, 0.5 * style_acc_current):
                acc_lab[i] = 2
            elif interp(series_vspd[i], vspd_interp, 0.15 * style_acc_current) < series_xacc[i] < \
                    interp(series_vspd[i], vspd_interp, 0.3 * style_acc_current):
                acc_lab[i] = 1
        acc_times = self.find_3_level_times(acc_lab)
        return acc_times

    def sudden_steering(self, series, car_type='default'):
        style_steer = {'ZS11': [0.35, 0.2, 0.1],
                       'AS24': [0.35, 0.2, 0.15],
                       'default': [0.35, 0.2, 0.15]}
        l = len(series)
        str_lab = zeros(l)
        for i in range(l):
            if series[i] > style_steer[car_type][0]:
                str_lab[i] = 3
            elif style_steer[car_type][0] > series[i] > style_steer[car_type][1]:
                str_lab[i] = 2
            elif style_steer[car_type][1] > series[i] > style_steer[car_type][2]:
                str_lab[i] = 1
        str_times = self.find_3_level_times(str_lab)
        return str_times

    @staticmethod
    def brake_skill(series_xacc, v_spd):
        l = len(v_spd)
        vel_spd_lab = zeros(l)
        brake_skill_cal = []
        for i in range(l):
            if not v_spd[i]:
                vel_spd_lab[i] = 1

        for i in range(15, l - 1, 1):  # 防止max_acc出错
            if not vel_spd_lab[i] and vel_spd_lab[i + 1]:
                max_acc = max(series_xacc[i - 15:i])
                I = series_xacc[i - 15:i].index(max_acc)
                # I = np.where(series_xacc[i - 15:i] == max_acc)
                start_ = I + i - 15 - 1
                for j in range(start_, start_ + 5, 1):
                    if series_xacc[j] > series_xacc[j + 1] and series_xacc[j + 1] <= series_xacc[j + 2]:
                        min_acc = series_xacc[j + 1]
                        brake_skill_cal = concatenate((brake_skill_cal, [max_acc - min_acc]))
                        break

        return array(brake_skill_cal).mean()

    @staticmethod
    def tip_in(series_acc_ped, frequency=10):
        l = len(series_acc_ped)
        delta_acc_ped = []
        result_array = []
        for i in range(1, l, 1):
            delta_acc_ped.append(series_acc_ped[i] - series_acc_ped[i - 1])
        delta_accped = array(delta_acc_ped) * frequency
        tip_in_times = 0
        for i in range(1, l - 5, 1):
            sum_of_acc = 0
            if delta_acc_ped[i] > 0 and delta_acc_ped[i - 1] <= 0:
                sum_of_acc = delta_accped[i]
                j = i + 1
                while delta_acc_ped[j] > 0 and delta_acc_ped[j + 3] >= 0 and j < l - 5:
                    sum_of_acc = sum_of_acc + delta_acc_ped[j]
                    j = j + 1
                if sum_of_acc / (j - i) > 80 and series_acc_ped[j - 1] > 10:
                    tip_in_times = tip_in_times + 1
                    result_array.append([series_acc_ped[i], series_acc_ped[j], (j - i) / frequency])
        return tip_in_times, result_array

    @staticmethod
    def soc_balance_time_ratio(series_soc_df, balance_position=33):
        time_ratio = series_soc_df[series_soc_df < balance_position].shape[0] / series_soc_df.shape[0]
        return time_ratio

    @staticmethod
    def spd_distribution_array(vspd, input_series, cal_type='sum'):
        output_array = []
        for i in range(1, 13, 1):
            if cal_type == 'sum':
                output_array.append(sum(input_series[(vspd < (i * 10)) & (vspd > ((i - 1) * 10))]))
            elif cal_type == 'avg':
                output_array.append(input_series[(vspd < (i * 10)) & (vspd > ((i - 1) * 10))].mean())
            elif cal_type == 'weight':
                output_array.append(input_series[(vspd < (i * 10)) & (vspd > ((i - 1) * 10))].size / input_series.size)

        return output_array

    @staticmethod
    def start_times_cal(vspd):
        start_times = 0
        l = vspd.shape[0]
        for i in range(1, l - 1):
            if vspd.iloc[i] == 0 and vspd.iloc[i + 1] > 0:
                start_times = start_times + 1
        return start_times

    # ------------------------   API

    def gps_fix_wgs84(self, lgt, lat):
        lgt = int(lgt / 100) + (lgt / 100 - int(lgt / 100)) * 100 / 60
        lat = int(lat / 100) + (lat / 100 - int(lat / 100)) * 100 / 60
        return self.wgs84_to_gcj02(lgt, lat)

    def wgs84_to_gcj02(self, lng, lat):
        if self.out_of_china(lng, lat):
            return [lng, lat]
        else:
            dlat = self.transform_lat(lng - 105.0, lat - 35.0)
            dlng = self.transform_lng(lng - 105.0, lat - 35.0)
            radlat = lat / 180.0 * self.PI
            magic = sin(radlat)
            magic = 1 - self.ee * magic * magic
            sqrtmagic = sqrt(magic)
            dlat = (dlat * 180.0) / ((self.a * (1 - self.ee)) / (magic * sqrtmagic) * self.PI)
            dlng = (dlng * 180.0) / (self.a / sqrtmagic * cos(radlat) * self.PI)
            mglat = lat + dlat
            mglng = lng + dlng
            return [mglng, mglat]

    def transform_lat(self, lng, lat):
        ret = -100.0 + 2.0 * lng + 3.0 * lat + 0.2 * lat * lat + 0.1 * lng * lat + 0.2 * sqrt(abs(lng))
        ret = ret + (20.0 * sin(6.0 * lng * self.PI) + 20.0 * sin(2.0 * lng * self.PI)) * 2.0 / 3.0
        ret = ret + (20.0 * sin(lat * self.PI) + 40.0 * sin(lat / 3.0 * self.PI)) * 2.0 / 3.0
        ret = ret + (160.0 * sin(lat / 12.0 * self.PI) + 320 * sin(lat * self.PI / 30.0)) * 2.0 / 3.0
        return ret

    def transform_lng(self, lng, lat):
        ret = 300.0 + lng + 2.0 * lat + 0.1 * lng * lng + 0.1 * lng * lat + 0.1 * sqrt(abs(lng))
        ret = ret + (20.0 * sin(6.0 * lng * self.PI) + 20.0 * sin(2.0 * lng * self.PI)) * 2.0 / 3.0
        ret = ret + (20.0 * sin(lng * self.PI) + 40.0 * sin(lng / 3.0 * self.PI)) * 2.0 / 3.0
        ret = ret + (150.0 * sin(lng / 12.0 * self.PI) + 300.0 * sin(lng / 30.0 * self.PI)) * 2.0 / 3.0
        return ret

    @staticmethod
    def out_of_china(lng, lat):
        return (lng < 72.004 or lng > 137.8347) or ((lat < 0.8293 or lat > 55.8271) or False)

    def gps_to_road_information(self, csvdata, i):
        csvdata = csvdata[['Longitude', 'Latitude', 'Direction', 'Date', 'Time', 'Veh_Spd_NonDrvn']]
        # csvdata = csvdata[csvdata['Longitude'] > 0]
        GPS = csvdata[['Longitude', 'Latitude']]
        Time = csvdata[['Date', 'Time']]
        Dirction = csvdata['Direction']
        Vspd = csvdata['Veh_Spd_NonDrvn']

        # # 另一种写法
        # server = 'http://restapi.amap.com/v3/autograsp'
        # parameters = {'key': '42f165dabcfcb28c1c0290058adee399', 'carid': 'e399123456',
        #               'locations': '116.496167,39.917066|116.496149,39.917205|116.496149,39.917326',
        #               'time': '1502242820,1502242823,1502242830', 'direction': '1,1,2',
        #               'speed': '1,1,2'}
        # r = requests.get(server, params=parameters)
        # print(r.url)

        # 每次取点
        base = 'http://restapi.amap.com/v3/autograsp?key=42f165dabcfcb28c1c0290058adee399&carid=e399123456'
        Gps1 = self.gps_fix_wgs84(GPS.iloc[i][0], GPS.iloc[i][1])
        Gps2 = self.gps_fix_wgs84(GPS.iloc[i + 50][0], GPS.iloc[i + 50][1])
        Gps3 = self.gps_fix_wgs84(GPS.iloc[i + 100][0], GPS.iloc[i + 100][1])
        location = '&locations=' + str(Gps1[0]) + ',' + str(Gps1[1]) + '|' + str(Gps2[0]) + ',' + str(
            Gps2[1]) + '|' + str(
            Gps3[0]) + ',' + str(Gps3[1])
        Time1 = self.utc_mktime(Time.iloc[i][0], Time.iloc[i][1])
        Time2 = self.utc_mktime(Time.iloc[i + 50][0], Time.iloc[i + 50][1])
        Time3 = self.utc_mktime(Time.iloc[i + 100][0], Time.iloc[i + 100][1])
        time = '&time=' + str(Time1) + ',' + str(Time2) + ',' + str(Time3)
        direction = '&direction=' + str(Dirction.iloc[i]) + ',' + str(Dirction.iloc[i + 50]) + ',' + str(
            Dirction.iloc[i + 100])
        speed = '&speed=' + str(Vspd.iloc[i]) + ',' + str(Vspd.iloc[i + 50]) + ',' + str(Vspd.iloc[i + 100])
        url = base + location + time + direction + speed
        r = get(url)
        answer = r.json()
        if answer['infocode'] == '10000':  # 请求成功
            if answer['roads'] != []:
                road = answer['roads'][0]['roadname']
                maxspeed = answer['roads'][0]['maxspeed']
                level = answer['roads'][0]['roadlevel']
            else:
                road = 'Unknown'
                maxspeed = '60'
                level = 'Unknown'
        else:
            road = 'GPS_error'
            maxspeed = '120'
            level = 'GPS_error'
        return road, maxspeed, level


if __name__ == '__main__':
    a = SystemGain('./IP31_L16UOV055_Ride_SyGa_20160225_SL.csv')
    a.sg_main()
    plt.show()

    print('Finish!')
